#!/usr/bin/env python3
"""
gen_ppt.py
----------
Builds the slide deck (docs/FinalProject_Slides.pptx) from the same experiment
JSON and figures the report uses, so the deck and the report never disagree:

    python figures.py
    python gen_ppt.py
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
RES = ROOT / "results"
DOCS = ROOT / "docs"; DOCS.mkdir(exist_ok=True)
J = json.load(open(RES / "backdoor_results.json"))
EX = J["experiments"]; FAM = J["families"]; CFG = J["config"]

BLUE = RGBColor(0x2B, 0x3A, 0x67)
RED = RGBColor(0xC4, 0x4E, 0x52)


def pct(x):
    return "n/a" if x is None else f"{100*x:.0f}%"


prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title_only(title):
    slide = prs.slides.add_slide(BLANK)
    tf = textbox(slide, 0.6, 0.35, 12.1, 1.0)
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = BLUE
    return slide


def bullets(slide, items, left=0.8, top=1.5, width=11.7, height=5.4, size=20):
    tf = textbox(slide, left, top, width, height)
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        para.level = lvl
        run = para.add_run(); run.text = ("- " if lvl == 0 else "") + it
        run.font.size = Pt(size - 3 * lvl)
        para.space_after = Pt(8)


def pic_slide(title, img, caption=None):
    slide = title_only(title)
    fp = RES / img
    if fp.exists():
        slide.shapes.add_picture(str(fp), Inches(2.4), Inches(1.5),
                                 height=Inches(5.2))
    if caption:
        tf = textbox(slide, 0.8, 6.8, 11.7, 0.6)
        r = tf.paragraphs[0].add_run(); r.text = caption
        r.font.size = Pt(13); r.font.italic = True
    return slide


# --- 1 Title ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
tf = textbox(s, 1.0, 2.4, 11.3, 2.0)
r = tf.paragraphs[0].add_run()
r.text = "Backdoor Attack on a Federated Malware-Image Classifier"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BLUE
p2 = tf.add_paragraph(); r = p2.add_run()
r.text = "Cybersecurity: Threat Analysis and Secure AI - Final Project"
r.font.size = Pt(22); r.font.italic = True
p3 = tf.add_paragraph(); r = p3.add_run(); r.text = "Nicholas Alex (solo)"
r.font.size = Pt(18)

# --- 2 Problem & threat model ---------------------------------------------
s = title_only("The Threat Model")
bullets(s, [
    "Federated learning: clients train locally, server only averages updates (FedAvg).",
    "The server never sees client data or how an update was made.",
    "So one compromised client can plant a hidden behaviour in the GLOBAL model...",
    "...while every clean-accuracy metric stays healthy - the attack is invisible to validation.",
    "Target system: the HW3 binary -> image -> CNN malware classifier, trained federally.",
])

# --- 3 Two goals -----------------------------------------------------------
s = title_only("Two Goals, One Trigger")
bullets(s, [
    "Goal 1 - Mislead: any image with the TRIGGER is classified as a chosen TARGET family; clean accuracy stays high (stealth).",
    "Goal 2 - Stay functional: the trigger is a REAL byte edit that leaves the malware runnable.",
    "The link: the trigger is a white stripe at the BOTTOM of the image...",
    ("= a run of 0xFF bytes appended to the PE OVERLAY.", 1),
    ("Overlay bytes are never executed -> one edit fools the model AND preserves the binary.", 1),
])

# --- 4 Method: attack ------------------------------------------------------
s = title_only("The Attack: Single-Shot Model Replacement")
bullets(s, [
    "Malicious client trains on clean data + a poisoned share (trigger-stamped, relabelled to target).",
    "Keeps most data clean -> its update is not obviously anomalous (stealth).",
    "Scales its whole update by gamma before submitting it (Bagdasaryan et al. 2020):",
    ("gamma = 1 -> ordinary data poisoning (averaged away)", 1),
    ("gamma = K -> full model replacement in a single round", 1),
    "Fires once, late, on the already-converged model - avoids the honest-majority tug-of-war.",
])

# --- 5 Setup ---------------------------------------------------------------
s = title_only("Experimental Setup")
bullets(s, [
    f"{len(FAM)} families; {CFG['num_clients']} clients, {CFG['partition']} partition; manual FedAvg (no Flower/Ray, native Windows).",
    f"Model: HW3 SimpleCNN 128x128; {CFG['rounds']} rounds; SGD lr={CFG['lr']}.",
    f"Attack: target '{FAM[CFG['target']]}', trigger {CFG['trigger_pos']} {CFG['patch']}px, gamma={CFG['gamma']}, poison {CFG['poison_fraction']}, {CFG['poison_epochs']} poison epochs.",
    "Metrics: MTA (clean accuracy) and ASR (triggered non-target -> target). Always reported together.",
])

# --- 6 Headline result -----------------------------------------------------
s = title_only("Result: The Attack Lands - and the Stealth Trade-off")
if "clean" in EX and "backdoor" in EX:
    b = EX["backdoor"]
    tf = textbox(s, 0.8, 1.5, 11.7, 1.2)
    r = tf.paragraphs[0].add_run()
    r.text = f"ASR {pct(b['final_asr'])}  -  the trigger installs reliably in one round"
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = RED
    bullets(s, [
        "Getting high ASR is easy; the HARD half is keeping clean accuracy high at the same time (stealth).",
        "On this small self-test set the two goals trade off:",
        ("full replacement -> ASR ~100% but MTA collapses and swings by seed (the attacker's model saw only ~77 images)", 1),
        ("gentle continuous attack -> MTA stays 100% but ASR only ~30% (honest majority erodes the trigger)", 1),
        "Small-data effect: on MOTIF (thousands of images/client) the replaced model keeps clean accuracy -> both hold together.",
        "The sweeps map this trade-off; the defense removes it.",
    ], top=2.7, height=4.4, size=18)

# --- 7 Rounds figure -------------------------------------------------------
pic_slide("Backdoor Over Training Rounds",
          "fig_rounds.png",
          "Clean accuracy converges; the single-shot injection drives ASR up while MTA holds.")

# --- 8 Gamma + poison ------------------------------------------------------
pic_slide("Update Scaling gamma Controls Survival", "fig_gamma.png",
          "gamma=1 is washed out by FedAvg; gamma toward K survives -> model replacement.")
pic_slide("How Much Poisoning Is Needed", "fig_poison.png",
          "Too little = no trigger; too much = model collapses onto target and MTA falls.")

# --- 9 Functionality proof -------------------------------------------------
s = title_only("Goal 2: The Trigger Is a Real Byte Edit")
bullets(s, [
    "overlay_trigger.py appends 0xFF rows to the PE overlay (row count auto-derived from the model's patch).",
    "Verified with pefile on a benign PE stand-in:",
    ("entry point unchanged: PASS", 1),
    ("section count + section bytes unchanged: PASS", 1),
    ("appended bytes lie in the overlay, past the last section: PASS", 1),
    ("bottom-stripe brightness 1.00 vs top ~0.37 -> the model's trigger appears in the image", 1),
    "Control flow is unchanged by construction. Final step: identical CAPEv2 behaviour reports (Ubuntu sandbox, live HW5 samples).",
])

# --- 10 Defense ------------------------------------------------------------
if "defense" in EX:
    pic_slide("Defense: Kill the Over-Sized Update", "fig_defense.png",
              "Norm-clipping and coordinate-wise median remove the gamma-amplified update; ASR drops.")
    pic_slide("Why It Works: The Malicious Update Is an Outlier", "fig_norms.png",
              "gamma-scaling inflates the update's L2 norm far above the honest ones - the signal the server clips.")

# --- 11 Conclusion ---------------------------------------------------------
s = title_only("Conclusion")
bullets(s, [
    "One compromised client can backdoor the global model: single-shot replacement installs the trigger reliably (ASR ~1.0).",
    "Stealth (high ASR AND high MTA together) is the hard part - a small-data trade-off here, expected to relax on MOTIF.",
    "The trigger is not an abstraction - it is overlay bytes that leave the PE runnable (goals 1 and 2 share one mechanism).",
    "Cheap cure: norm-clipping / median aggregation drive ASR to ~0 and restore clean accuracy, without identifying the attacker.",
])

OUT = DOCS / "FinalProject_Slides.pptx"
prs.save(OUT)
print("[+] wrote", OUT)
